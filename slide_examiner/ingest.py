from __future__ import annotations

import json
import re
import zipfile
from html.parser import HTMLParser
from pathlib import Path
from typing import Any

from .schemas import BBox, Deck, Element, Slide


def load_slide_json(path: str | Path) -> Slide:
    return Slide.from_mapping(json.loads(Path(path).read_text(encoding="utf-8")))


def load_deck_json(path: str | Path) -> Deck:
    value = json.loads(Path(path).read_text(encoding="utf-8"))
    if "slides" in value:
        return Deck.from_mapping(value)
    return Deck(deck_id=Path(path).stem, slides=(Slide.from_mapping(value),))


def save_slide_json(slide: Slide, path: str | Path) -> None:
    output = Path(path)
    output.parent.mkdir(parents=True, exist_ok=True)
    output.write_text(json.dumps(slide.to_dict(), ensure_ascii=False, indent=2), encoding="utf-8")


def save_deck_json(deck: Deck, path: str | Path) -> None:
    output = Path(path)
    output.parent.mkdir(parents=True, exist_ok=True)
    output.write_text(json.dumps(deck.to_dict(), ensure_ascii=False, indent=2), encoding="utf-8")


def parse_annotated_html(path: str | Path, *, slide_id: str | None = None, width: int = 1920, height: int = 1080) -> Slide:
    parser = _AnnotatedHTMLParser(slide_id=slide_id or Path(path).stem, width=width, height=height)
    parser.feed(Path(path).read_text(encoding="utf-8"))
    return parser.slide()


def extract_pptx_text_outline(path: str | Path) -> Deck:
    """Extract a lightweight text-only IR from PPTX without optional dependencies.

    This is not a replacement for python-pptx geometry extraction. It gives the
    pipeline a deterministic fallback and makes missing geometry explicit.
    """

    pptx_path = Path(path)
    slides: list[Slide] = []
    with zipfile.ZipFile(pptx_path) as archive:
        slide_names = sorted(
            name for name in archive.namelist() if re.fullmatch(r"ppt/slides/slide\d+\.xml", name)
        )
        for index, name in enumerate(slide_names, start=1):
            xml = archive.read(name).decode("utf-8", errors="ignore")
            texts = [re.sub(r"<.*?>", "", match) for match in re.findall(r"<a:t>(.*?)</a:t>", xml)]
            elements = tuple(
                Element(
                    element_id=f"s{index}_text{item_index}",
                    type="text",
                    bbox=BBox(96, 96 + item_index * 72, 1728, 56),
                    text=text,
                    metadata={"source": "pptx_text_outline", "geometry_confidence": "fallback"},
                )
                for item_index, text in enumerate(texts)
            )
            slides.append(Slide(slide_id=f"{pptx_path.stem}_{index}", elements=elements))
    return Deck(deck_id=pptx_path.stem, slides=tuple(slides), metadata={"source_path": str(pptx_path)})


def extract_pptx_geometry(path: str | Path, *, fallback: bool = True) -> Deck:
    try:
        from pptx import Presentation
    except ImportError:
        if fallback:
            return extract_pptx_text_outline(path)
        raise

    pptx_path = Path(path)
    try:
        presentation = Presentation(str(pptx_path))
    except Exception:
        if fallback:
            return extract_pptx_text_outline(path)
        raise

    width_emu = float(presentation.slide_width)
    height_emu = float(presentation.slide_height)
    slides: list[Slide] = []
    for slide_index, ppt_slide in enumerate(presentation.slides, start=1):
        elements: list[Element] = []
        for shape_index, shape in enumerate(ppt_slide.shapes):
            bbox = BBox(
                x=float(shape.left) / width_emu * 1920,
                y=float(shape.top) / height_emu * 1080,
                width=float(shape.width) / width_emu * 1920,
                height=float(shape.height) / height_emu * 1080,
            )
            text = getattr(shape, "text", "") or ""
            style = _shape_style(shape)
            metadata = {
                "source": "python-pptx",
                "shape_id": getattr(shape, "shape_id", None),
                "shape_name": getattr(shape, "name", None),
                "geometry_confidence": "extracted",
            }
            elements.append(
                Element(
                    element_id=f"s{slide_index}_shape{shape_index}",
                    type="text" if text else "shape",
                    bbox=bbox,
                    text=text,
                    style=style,
                    z=shape_index,
                    metadata={key: value for key, value in metadata.items() if value is not None},
                )
            )
        slides.append(Slide(slide_id=f"{pptx_path.stem}_{slide_index}", elements=tuple(elements)))
    return Deck(deck_id=pptx_path.stem, slides=tuple(slides), metadata={"source_path": str(pptx_path), "source": "python-pptx"})


def slide_caption(slide: Slide) -> str:
    parts = []
    for element in slide.elements:
        descriptor = f"{element.type} {element.element_id} at ({element.bbox.x:.0f},{element.bbox.y:.0f},{element.bbox.width:.0f},{element.bbox.height:.0f})"
        if element.text:
            descriptor += f": {element.text[:120]}"
        parts.append(descriptor)
    return " | ".join(parts)


def deck_caption(deck: Deck) -> str:
    return " || ".join(f"Slide {index + 1}: {slide_caption(slide)}" for index, slide in enumerate(deck.slides))


class _AnnotatedHTMLParser(HTMLParser):
    def __init__(self, *, slide_id: str, width: int, height: int) -> None:
        super().__init__()
        self.slide_id = slide_id
        self.width = width
        self.height = height
        self._stack: list[dict[str, Any]] = []
        self._elements: list[Element] = []

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        attr = {key: value or "" for key, value in attrs}
        element_id = attr.get("data-element-id") or attr.get("id")
        if not element_id or "data-x" not in attr:
            self._stack.append({"capture": False})
            return
        style = _parse_style(attr.get("style", ""))
        bbox = BBox(
            float(attr.get("data-x", style.get("left", 0))),
            float(attr.get("data-y", style.get("top", 0))),
            float(attr.get("data-width", style.get("width", 0))),
            float(attr.get("data-height", style.get("height", 0))),
        )
        metadata = {
            key.removeprefix("data-meta-"): value
            for key, value in attr.items()
            if key.startswith("data-meta-")
        }
        self._stack.append(
            {
                "capture": True,
                "element_id": element_id,
                "type": attr.get("data-type", tag),
                "bbox": bbox,
                "style": {
                    "font_size_pt": _maybe_float(attr.get("data-font-size")),
                    "color": attr.get("data-color") or style.get("color"),
                    "fill_color": attr.get("data-fill-color") or style.get("background-color"),
                },
                "metadata": metadata,
                "text": [],
            }
        )

    def handle_data(self, data: str) -> None:
        if self._stack and self._stack[-1].get("capture"):
            self._stack[-1]["text"].append(data)

    def handle_endtag(self, tag: str) -> None:
        if not self._stack:
            return
        frame = self._stack.pop()
        if not frame.get("capture"):
            return
        style = {key: value for key, value in frame["style"].items() if value is not None}
        self._elements.append(
            Element(
                element_id=frame["element_id"],
                type=frame["type"],
                bbox=frame["bbox"],
                text="".join(frame["text"]).strip(),
                style=style,
                metadata=frame["metadata"],
            )
        )

    def slide(self) -> Slide:
        return Slide(slide_id=self.slide_id, elements=tuple(self._elements), width=self.width, height=self.height)


def _parse_style(style: str) -> dict[str, Any]:
    parsed: dict[str, Any] = {}
    for item in style.split(";"):
        if ":" not in item:
            continue
        key, value = item.split(":", 1)
        key = key.strip()
        value = value.strip()
        if key in {"left", "top", "width", "height"}:
            parsed[key] = _css_px(value)
        else:
            parsed[key] = value
    return parsed


def _css_px(value: str) -> float:
    return float(value.replace("px", "").strip())


def _maybe_float(value: str | None) -> float | None:
    if value in {None, ""}:
        return None
    return float(value)


def _shape_style(shape) -> dict[str, Any]:
    style: dict[str, Any] = {}
    if getattr(shape, "has_text_frame", False) and shape.text_frame.paragraphs:
        paragraph = shape.text_frame.paragraphs[0]
        if paragraph.runs:
            run = paragraph.runs[0]
            if run.font.size is not None:
                style["font_size_pt"] = float(run.font.size.pt)
            rgb = _font_rgb(run)
            if rgb is not None:
                style["color"] = rgb
    return style


def _font_rgb(run) -> str | None:
    try:
        rgb = run.font.color.rgb
    except Exception:
        return None
    if rgb is None:
        return None
    return f"#{rgb}"
