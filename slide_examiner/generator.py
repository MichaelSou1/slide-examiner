"""Part 3 slide-generator agent (the *optimized system* / rollout engine).

The agent turns a task brief into a structured content JSON and then into a
rendered deck (image + structure), steered entirely by the four editable
``PromptModules`` (the optimizer's edit surface). The structural layout is fixed
so that *skill quality*, not code, drives changes in the generated deck.

``deck_from_content_json`` (the content-JSON -> IR step) is shared with Parts 1/2
and is extended here only additively (figures + key_terms); the title+bullets
skeleton is unchanged.
"""
from __future__ import annotations

import json
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Callable

from .adapters import parse_examiner_json
from .api_config import build_completion
from .render import render_slide_html_file, render_slide_multi_resolution
from .schemas import BBox, Deck, Element, Slide
from .skill_doc import (  # re-exported for backwards-compatible imports
    DEFAULT_PROMPT_MODULES,
    PromptModules,
    render_skill_doc,
)

__all__ = [
    "PromptModules",
    "DEFAULT_PROMPT_MODULES",
    "deck_from_content_json",
    "load_content_json",
    "write_deck_html",
    "export_deck_pptx",
    "GeneratorConfig",
    "GeneratedArtifact",
    "generate_deck",
    "build_generation_messages",
]


# --------------------------------------------------------------------------- #
# content JSON  ->  Deck IR  (shared with Parts 1/2; figures/key_terms added)
# --------------------------------------------------------------------------- #
def deck_from_content_json(content: dict[str, Any], *, prompt_modules: PromptModules | None = None) -> Deck:
    modules = prompt_modules or PromptModules()
    deck_id = str(content.get("deck_id", "generated_deck"))
    slides: list[Slide] = []
    for index, page in enumerate(content.get("slides", []), start=1):
        title = str(page.get("title", f"Slide {index}"))
        bullets = [str(item) for item in page.get("bullets", [])]
        section = page.get("section")
        elements = [
            Element(
                element_id=f"s{index}_title",
                type="title",
                bbox=BBox(96, 72, 1728, 72),
                text=title,
                style={"font_size_pt": 34, "color": "#111111"},
                metadata={"role": "title", "text_level": "title"},
            )
        ]
        next_y = 180
        for bullet_index, bullet in enumerate(bullets):
            elements.append(
                Element(
                    element_id=f"s{index}_body{bullet_index}",
                    type="text",
                    bbox=BBox(144, 180 + bullet_index * 92, 1512, 72),
                    text=bullet,
                    style={"font_size_pt": 22, "color": "#222222"},
                    metadata={"role": "body", "text_level": "body"},
                )
            )
            next_y = 180 + (bullet_index + 1) * 92

        # additive: figure elements (pixel-visible claim text, feeds S6)
        figures = page.get("figures") or []
        for fig_index, figure in enumerate(figures):
            if not isinstance(figure, dict):
                continue
            kind = str(figure.get("kind", "trend_up"))
            claim = str(figure.get("claim", "")).strip()
            if not claim:
                continue
            elements.append(
                Element(
                    element_id=f"s{index}_fig{fig_index}",
                    type="figure",
                    bbox=BBox(144, next_y + 24 + fig_index * 96, 1512, 80),
                    text=claim,
                    style={"font_size_pt": 20, "color": "#0a5", "fill_color": "#eef7f0"},
                    metadata={"role": "figure", "kind": kind, "diagram_claim": claim},
                )
            )

        key_terms = [str(t) for t in (page.get("key_terms") or []) if str(t).strip()]
        slides.append(
            Slide(
                slide_id=f"{deck_id}_{index}",
                elements=tuple(elements),
                metadata={
                    "section": section,
                    "page_type": page.get("page_type"),
                    "key_terms": key_terms,
                    "prompt_modules": modules.to_dict(),
                },
            )
        )
    return Deck(
        deck_id=deck_id,
        slides=tuple(slides),
        metadata={
            "scenario": content.get("scenario"),
            "required_sections": content.get("required_sections", []),
            "project_glossary": content.get("key_terms", []),
            "prompt_modules": modules.to_dict(),
        },
    )


def load_content_json(path: str | Path) -> dict[str, Any]:
    return json.loads(Path(path).read_text(encoding="utf-8"))


def write_deck_html(deck: Deck, output_dir: str | Path) -> list[Path]:
    base = Path(output_dir)
    base.mkdir(parents=True, exist_ok=True)
    paths = []
    for index, slide in enumerate(deck.slides, start=1):
        paths.append(render_slide_html_file(slide, base / f"{index:03d}_{slide.slide_id}.html"))
    return paths


def export_deck_pptx(deck: Deck, output_path: str | Path) -> Path:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError as exc:
        raise RuntimeError("python-pptx is required for PPTX export.") from exc

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for slide_ir in deck.slides:
        ppt_slide = prs.slides.add_slide(blank)
        for element in slide_ir.elements:
            left = Inches(element.bbox.x / slide_ir.width * 13.333)
            top = Inches(element.bbox.y / slide_ir.height * 7.5)
            width = Inches(element.bbox.width / slide_ir.width * 13.333)
            height = Inches(element.bbox.height / slide_ir.height * 7.5)
            box = ppt_slide.shapes.add_textbox(left, top, width, height)
            box.text = element.text
            for paragraph in box.text_frame.paragraphs:
                for run in paragraph.runs:
                    if "font_size_pt" in element.style:
                        run.font.size = Pt(float(element.style["font_size_pt"]))
    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)
    return output


# --------------------------------------------------------------------------- #
# Generator agent: brief  ->  content JSON  ->  Deck IR  ->  rendered deck
# --------------------------------------------------------------------------- #
@dataclass(frozen=True)
class GeneratorConfig:
    """Config for the local text-LLM that drives content synthesis.

    Defaults target a locally-served Qwen3.6-27B (vLLM, OpenAI-compatible). The
    generator weights are FROZEN; only the skill modules change across rollouts.
    """

    model: str = "qwen3.6-27b"
    base_url: str | None = "http://127.0.0.1:8200/v1"
    api_key_env: str = "OPENAI_API_KEY"
    api_style: str = "chat"  # "chat" | "responses"
    max_tokens: int = 4096
    temperature: float = 0.0
    renderer: str = "html"  # "html" | "pptx"
    long_edge: int = 1024
    max_slides: int = 15


@dataclass(frozen=True)
class GeneratedArtifact:
    deck: Deck
    content_json: dict[str, Any]
    page_image_paths: list[Path]
    render_specs: list[dict[str, Any]]
    prompt_modules: PromptModules
    raw_completion: str
    out_dir: Path
    degenerate: bool = False

    def to_dict(self) -> dict[str, Any]:
        return {
            "deck_id": self.deck.deck_id,
            "n_slides": len(self.deck.slides),
            "page_image_paths": [str(p) for p in self.page_image_paths],
            "degenerate": self.degenerate,
            "out_dir": str(self.out_dir),
        }


_SYSTEM_PREAMBLE = (
    "You are a slide-deck generator. Given a task brief and four skill modules, "
    "produce ONE presentation as a single JSON object. Follow the skill modules "
    "exactly. Output ONLY the JSON object, no prose, no markdown fences."
)

_CONTENT_SCHEMA_HINT = (
    "JSON schema:\n"
    "{\n"
    '  "deck_id": str,\n'
    '  "scenario": "launch" | "client_intro" | "full_proposal",\n'
    '  "required_sections": [str, ...],\n'
    '  "slides": [\n'
    "    {\n"
    '      "page_type": "title"|"agenda"|"content"|"comparison"|"closing",\n'
    '      "section": str,\n'
    '      "title": str,\n'
    '      "bullets": [str, ...],\n'
    '      "figures": [{"kind": "trend_up"|"trend_down"|"arch", "claim": str}],\n'
    '      "key_terms": [str, ...]\n'
    "    }\n"
    "  ]\n"
    "}"
)


def build_generation_messages(task: dict[str, Any], modules: PromptModules) -> list[dict[str, str]]:
    """Build the chat messages; the four skill modules are the editable surface."""

    brief = str(task.get("brief") or task.get("task_brief") or task.get("description") or "")
    deck_id = str(task.get("task_id") or task.get("id") or "generated_deck")
    required = task.get("required_sections")
    skill_block = render_skill_doc(modules)
    user = [
        f"TASK_ID: {deck_id}",
        f"BRIEF:\n{brief}",
    ]
    if required:
        user.append(f"REQUIRED_SECTIONS (must each appear): {', '.join(map(str, required))}")
    user.append("")
    user.append("SKILL MODULES (follow exactly):")
    user.append(skill_block)
    user.append(_CONTENT_SCHEMA_HINT)
    user.append(f'Set "deck_id" to "{deck_id}".')
    return [
        {"role": "system", "content": _SYSTEM_PREAMBLE},
        {"role": "user", "content": "\n".join(user)},
    ]


def _default_complete(config: GeneratorConfig) -> Callable[[list[dict[str, str]]], str]:
    """Text completion against an OpenAI-compatible endpoint (local vLLM or API)."""

    return build_completion(
        config.model,
        config.base_url,
        api_key_env=config.api_key_env,
        api_style=config.api_style,
        max_tokens=config.max_tokens,
        temperature=config.temperature,
    )


def _empty_deck(deck_id: str, modules: PromptModules) -> Deck:
    return Deck(deck_id=deck_id, slides=tuple(), metadata={"prompt_modules": modules.to_dict(), "degenerate": True})


def generate_deck(
    task: dict[str, Any],
    prompt_modules: PromptModules,
    config: GeneratorConfig | None = None,
    *,
    out_dir: str | Path,
    seed: int = 0,
    complete: Callable[[list[dict[str, str]]], str] | None = None,
    render: bool = True,
    revise: str | None = None,
) -> GeneratedArtifact:
    """Run one generation rollout: brief -> content JSON -> deck IR -> render.

    A degenerate (unparseable) completion yields an empty deck with
    ``degenerate=True`` so the feedback scorer can penalize it — the loop never
    crashes on a bad skill edit. ``complete`` may be injected for tests; when
    omitted a local vLLM client is used. ``seed`` varies the prompt (and the
    sampler when supported) so rollouts are reproducible per seed. ``revise``
    (self-refine) appends a user turn carrying the previous deck + critique so the
    generator REVISES rather than generates from scratch.
    """

    config = config or GeneratorConfig()
    deck_id = str(task.get("task_id") or task.get("id") or "generated_deck")
    out = Path(out_dir)
    out.mkdir(parents=True, exist_ok=True)

    messages = build_generation_messages(task, prompt_modules)
    if revise:
        messages = [*messages, {"role": "user", "content": revise}]
    if seed:
        messages = [*messages, {"role": "system", "content": f"(generation seed: {seed})"}]
    runner = complete or _default_complete(config)

    raw = ""
    content: dict[str, Any] | None = None
    for attempt in range(2):
        try:
            raw = runner(messages if attempt == 0 else [*messages, {"role": "user", "content": "OUTPUT VALID JSON ONLY."}])
            parsed = parse_examiner_json(raw)
            if isinstance(parsed, dict) and parsed.get("slides"):
                content = parsed
                break
        except Exception:
            content = None
    if content is None:
        return GeneratedArtifact(
            deck=_empty_deck(deck_id, prompt_modules),
            content_json={},
            page_image_paths=[],
            render_specs=[],
            prompt_modules=prompt_modules,
            raw_completion=raw,
            out_dir=out,
            degenerate=True,
        )

    content.setdefault("deck_id", deck_id)
    if config.max_slides and isinstance(content.get("slides"), list):
        content["slides"] = content["slides"][: config.max_slides]
    deck = deck_from_content_json(content, prompt_modules=prompt_modules)

    page_image_paths: list[Path] = []
    render_specs: list[dict[str, Any]] = []
    if render and deck.slides:
        for slide in deck.slides:
            artifacts = render_slide_multi_resolution(
                slide, out / slide.slide_id, long_edges=(config.long_edge,), filename="page"
            )
            page_image_paths.append(artifacts[0].image_path)
            render_specs.append(artifacts[0].render_spec)
        (out / "content.json").write_text(json.dumps(content, ensure_ascii=False, indent=2), encoding="utf-8")

    return GeneratedArtifact(
        deck=deck,
        content_json=content,
        page_image_paths=page_image_paths,
        render_specs=render_specs,
        prompt_modules=prompt_modules,
        raw_completion=raw,
        out_dir=out,
        degenerate=False,
    )
