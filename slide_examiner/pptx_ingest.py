"""Ingest a real .pptx into the slide-examiner Deck IR (Part 3 Hermes case-study).

This is intentionally lightweight: it reads top-level shape geometry + text via
python-pptx so a REAL agent-produced deck (e.g. Hermes pre-sales PPT) can be fed
to the same examiner contract (modality C: rendered image + element structure)
and the offline geometry linter. Group shapes contribute concatenated descendant
text at the group's bounding box (exact nested geometry is not needed — the
rendered image is the primary signal in modality C).
"""
from __future__ import annotations

from pathlib import Path
from typing import Any

from .schemas import BBox, Deck, Element, Slide


def _emu_to_px(v, scale: float) -> float:
    return float((v or 0)) * scale


def _shape_text(shape) -> str:
    """Text for a shape; recurse into groups to gather descendant text."""
    try:
        if shape.has_text_frame:
            return (shape.text_frame.text or "").strip()
    except Exception:
        pass
    # group shape -> concatenate descendant text
    try:
        if shape.shape_type is not None and int(shape.shape_type) == 6:  # GROUP
            parts = [_shape_text(c) for c in shape.shapes]
            return "\n".join(p for p in parts if p).strip()
    except Exception:
        pass
    return ""


def _shape_type(shape) -> str:
    """Map a python-pptx shape to the examiner ElementType vocabulary."""
    try:
        if shape.is_placeholder:
            pht = str(shape.placeholder_format.type).lower()
            if "title" in pht:
                return "subtitle" if "subtitle" in pht else "title"
            if "body" in pht or "object" in pht or "content" in pht:
                return "body"
    except Exception:
        pass
    try:
        st = int(shape.shape_type) if shape.shape_type is not None else -1
    except Exception:
        st = -1
    if st == 13:  # PICTURE
        return "image"
    if st == 19:  # TABLE
        return "table"
    if st == 3:   # CHART
        return "chart"
    try:
        if shape.has_text_frame and (shape.text_frame.text or "").strip():
            return "body"
    except Exception:
        pass
    return "shape"


def deck_from_pptx(path: str | Path, *, render_w: int = 2001, render_h: int = 1125,
                   scene: str = "full_proposal", deck_id: str | None = None) -> Deck:
    from pptx import Presentation  # local import (heavy)

    prs = Presentation(str(path))
    sw = prs.slide_width or 12192000
    sh = prs.slide_height or 6858000
    sx, sy = render_w / float(sw), render_h / float(sh)
    slides: list[Slide] = []
    for i, s in enumerate(prs.slides):
        els: list[Element] = []
        for j, shape in enumerate(s.shapes):
            text = _shape_text(shape)
            bbox = BBox(
                x=_emu_to_px(shape.left, sx), y=_emu_to_px(shape.top, sy),
                width=_emu_to_px(shape.width, sx), height=_emu_to_px(shape.height, sy),
            )
            els.append(Element(
                element_id=f"s{i+1:02d}_e{j:02d}", type=_shape_type(shape), bbox=bbox,
                text=text, style={}, z=j, placeholder_id=None, metadata={},  # always str (linter does len())
            ))
        slides.append(Slide(slide_id=f"s{i+1:02d}", elements=els, width=render_w,
                            height=render_h, metadata={"scene": scene}))
    return Deck(deck_id=deck_id or Path(path).stem, slides=slides, metadata={"scene": scene})


def placeholder_stats(deck: Deck, markers: tuple[str, ...] = ("添加", "点击此处", "请输入", "lorem", "your text", "add ")) -> dict[str, Any]:
    """Count unfilled-template-placeholder elements (a real Hermes failure mode).

    A verifiable, model-free signal: how many visible text elements still contain a
    template placeholder marker (e.g. 添加标题 / 添加内文). Used as the case-study DV
    alongside the examiner scores."""
    per_slide: list[int] = []
    total = 0
    for s in deck.slides:
        n = 0
        for e in s.elements:
            t = (e.text or "")
            if any(m in t.lower() if m.isascii() else m in t for m in markers):
                n += 1
        per_slide.append(n)
        total += n
    return {"total_placeholders": total, "per_slide": per_slide,
            "slides_with_placeholders": sum(1 for n in per_slide if n > 0), "n_slides": len(deck.slides)}
