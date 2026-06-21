"""Tests for pptx -> Deck IR ingest (Part 3 Hermes case-study)."""
import pytest

pptx = pytest.importorskip("pptx")

from slide_examiner.pptx_ingest import deck_from_pptx, placeholder_stats


def _make_pptx(path):
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    blank = prs.slide_layouts[6]
    # slide 1: a filled title + an unfilled placeholder marker
    s1 = prs.slides.add_slide(blank)
    tb = s1.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    tb.text_frame.text = "智能制造售前方案"
    tb2 = s1.shapes.add_textbox(Inches(1), Inches(3), Inches(5), Inches(1))
    tb2.text_frame.text = "添加标题"
    # slide 2: clean
    s2 = prs.slides.add_slide(blank)
    tb3 = s2.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    tb3.text_frame.text = "真实内容，无占位符"
    prs.save(str(path))


def test_deck_from_pptx_and_placeholder_stats(tmp_path):
    p = tmp_path / "d.pptx"
    _make_pptx(p)
    deck = deck_from_pptx(p, render_w=2001, render_h=1125)
    assert len(deck.slides) == 2
    # every element text is a str (linter does len()) and bboxes are in px
    for s in deck.slides:
        for e in s.elements:
            assert isinstance(e.text, str)
            assert e.bbox.width >= 0 and e.bbox.height >= 0
    ps = placeholder_stats(deck)
    assert ps["total_placeholders"] == 1
    assert ps["slides_with_placeholders"] == 1
    assert ps["n_slides"] == 2
