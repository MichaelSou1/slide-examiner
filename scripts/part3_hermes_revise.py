"""Part 3 Hermes case-study — one revision pass (examiner-critique-driven).

Takes the real deck's verifiable defect (unfilled template placeholders the
examiner flagged) and asks the GENERATOR (Hermes's own mimo-v2.5-pro, NOT qwen) to
fill them in, given deck context. Writes the replacements back into the .pptx via
python-pptx, re-renders, and re-counts placeholders to show the mess drops.

This is the lightweight "examiner critique -> generator revision -> measurable
improvement" loop, on a REAL agent deck. mimo-v2.5-pro = PART3_OPTIMIZER_* role.
"""
from __future__ import annotations

import argparse
import json
import re
import subprocess
import sys
from pathlib import Path

REPO = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(REPO))

from slide_examiner.api_config import build_completion, load_dotenv, resolve_role
from slide_examiner.pptx_ingest import deck_from_pptx, placeholder_stats

load_dotenv(REPO / ".env")
_GEN = resolve_role("OPTIMIZER", default_model="mimo-v2.5-pro")  # Hermes's configured mimo (not qwen)

MARKERS = ("添加", "点击此处", "请输入", "占位")


def _is_ph(text: str) -> bool:
    return any(m in (text or "") for m in MARKERS)


def _iter_text_shapes(shapes):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    for sh in shapes:
        try:
            if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
                yield from _iter_text_shapes(sh.shapes)
                continue
        except Exception:
            pass
        if getattr(sh, "has_text_frame", False):
            yield sh


def _set_text(shape, new_text: str) -> None:
    tf = shape.text_frame
    paras = tf.paragraphs
    p0 = paras[0]
    if p0.runs:
        p0.runs[0].text = new_text
        for r in p0.runs[1:]:
            r.text = ""
    else:
        p0.text = new_text
    for p in paras[1:]:
        for r in p.runs:
            r.text = ""


def _deck_brief(prs) -> str:
    """Best-effort deck topic from the first slide's filled text."""
    bits = []
    for sh in _iter_text_shapes(prs.slides[0].shapes):
        t = (sh.text_frame.text or "").strip()
        if t and not _is_ph(t):
            bits.append(t)
    return " / ".join(bits[:4]) or "盛原成 智能制造 / 数字化转型 售前方案"


def main() -> None:
    ap = argparse.ArgumentParser()
    ap.add_argument("--pptx", default=str(REPO / "data/part3/hermes_case/raw/deck.pptx"))
    ap.add_argument("--out-pptx", default=str(REPO / "data/part3/hermes_case/revised/deck.pptx"))
    ap.add_argument("--render", action="store_true", help="re-render the revised deck to PNGs")
    ap.add_argument("--out", default=str(REPO / "runs/probe/part3/hermes_revise.json"))
    args = ap.parse_args()

    from pptx import Presentation
    prs = Presentation(args.pptx)
    brief = _deck_brief(prs)
    complete = build_completion(_GEN["model"], _GEN["base_url"], api_key_env=_GEN["api_key_env"],
                               api_style=_GEN["api_style"], max_tokens=1200)

    before = placeholder_stats(deck_from_pptx(args.pptx))
    print(f"[revise] deck brief: {brief}")
    print(f"[revise] BEFORE: {before['total_placeholders']} placeholders / {before['slides_with_placeholders']} slides")

    filled = 0
    per_slide_log = []
    for si, slide in enumerate(prs.slides):
        shapes = list(_iter_text_shapes(slide.shapes))
        ph_shapes = [sh for sh in shapes if _is_ph(sh.text_frame.text)]
        if not ph_shapes:
            continue
        context = [sh.text_frame.text.strip() for sh in shapes if not _is_ph(sh.text_frame.text) and sh.text_frame.text.strip()]
        slots = [{"id": i, "placeholder": sh.text_frame.text.strip()[:40]} for i, sh in enumerate(ph_shapes)]
        prompt = (
            "你在修订一份真实的售前方案 PPT 的某一页。该页含未填写的模板占位符(如 添加标题/添加内文)。\n"
            f"全篇主题: {brief}\n"
            f"本页已有真实文字(上下文): {context}\n"
            f"待填占位符槽位(JSON): {json.dumps(slots, ensure_ascii=False)}\n"
            "请为每个槽位生成简洁、专业、与主题一致的中文实际内容(标题<=14字; 正文<=40字, 不要换行)。\n"
            '只输出 JSON: {"fills":[{"id":int,"text":str}, ...]}'
        )
        try:
            raw = complete([{"role": "user", "content": prompt}])
            m = re.search(r"\{.*\}", raw, re.DOTALL)
            fills = json.loads(m.group(0))["fills"] if m else []
        except Exception as exc:
            print(f"  [s{si+1:02d}] mimo failed: {type(exc).__name__}: {exc}")
            continue
        by_id = {int(f["id"]): str(f["text"]).strip() for f in fills if "id" in f and f.get("text")}
        slide_filled = 0
        for i, sh in enumerate(ph_shapes):
            if i in by_id and by_id[i]:
                _set_text(sh, by_id[i])
                slide_filled += 1
        filled += slide_filled
        per_slide_log.append({"slide": si + 1, "n_placeholders": len(ph_shapes), "n_filled": slide_filled})
        print(f"  [s{si+1:02d}] filled {slide_filled}/{len(ph_shapes)} placeholders", flush=True)

    Path(args.out_pptx).parent.mkdir(parents=True, exist_ok=True)
    prs.save(args.out_pptx)
    after = placeholder_stats(deck_from_pptx(args.out_pptx))
    print(f"[revise] AFTER:  {after['total_placeholders']} placeholders / {after['slides_with_placeholders']} slides "
          f"(filled {filled})")

    render_pngs = None
    if args.render:
        rdir = str(Path(args.out_pptx).parent)
        subprocess.run(["bash", str(REPO / "scripts/render_pptx_extracted.sh"), args.out_pptx, rdir, "150"],
                       check=False)
        render_pngs = len(list(Path(rdir).glob("slide-*.png")))

    out = {
        "deck": Path(args.pptx).name, "generator": _GEN["model"], "brief": brief,
        "before": before, "after": after, "n_filled": filled, "per_slide": per_slide_log,
        "rendered_pngs": render_pngs,
        "note": "Examiner-critique-driven revision: mimo-v2.5-pro (Hermes's model) fills the unfilled "
                "template placeholders the examiner flagged; placeholder count is the verifiable DV.",
    }
    Path(args.out).parent.mkdir(parents=True, exist_ok=True)
    Path(args.out).write_text(json.dumps(out, indent=2, ensure_ascii=False), encoding="utf-8")
    print(f"-> {args.out}")


if __name__ == "__main__":
    main()
