#!/usr/bin/env python
"""G7 render-containment-overflow *detector* for real .pptx decks (todo_0625 R5).

G7 = a text box whose *declared* geometry is legal (inside page margins, not
overlapping another box) but whose *rendered* content overflows the box. The
declared-geometry linter is blind to it by construction; the only way to see it
is to render faithfully and compare rendered text extent against the declared
box.

This runs the G7 containment check as a DETECTOR on un-injected real decks:

  declared box geometry   <- python-pptx (shape.left/top/width/height, EMU)
  rendered text extent    <- LibreOffice pptx->pdf, then PyMuPDF vector text bboxes
  faithful renderer       <- slide_examiner.render.render_pptx_to_pdf (soffice)

We deliberately do NOT use slide_examiner.render.slide_to_html: that path bakes
`overflow:visible` into every element, which would make the measurement circular.

Incidence is reported conditioned on *legal* text boxes only (the G7 definition),
so unfilled placeholders / declared-overflow (G1) do not inflate the rate.

Scope: text overflow (the card_height / unbreakable-text G7 variants). Image
bleed (image_objectfit) needs pixel methods and is out of scope here.
"""
from __future__ import annotations

import argparse
import json
import sys
import tempfile
from dataclasses import dataclass, field
from pathlib import Path

REPO = Path(__file__).resolve().parents[1]
if str(REPO) not in sys.path:
    sys.path.insert(0, str(REPO))

EMU_PER_PT = 12700.0  # 914400 EMU/in / 72 pt/in


def emu2pt(v) -> float:
    return float(v) / EMU_PER_PT if v is not None else 0.0


@dataclass
class Box:
    sid: str
    x0: float
    y0: float
    x1: float
    y1: float
    text: str
    auto_size: str | None = None
    word_wrap: bool | None = None

    @property
    def w(self) -> float:
        return self.x1 - self.x0

    @property
    def h(self) -> float:
        return self.y1 - self.y0


@dataclass
class Hit:
    sid: str
    box: tuple
    rendered_bottom: float
    rendered_right: float
    v_overflow_pt: float
    h_overflow_pt: float
    v_overflow_frac: float
    h_overflow_frac: float
    auto_size: str | None
    text_preview: str
    page_overflow_pt: float = 0.0  # rendered extent beyond the page edge
    collides: bool = False  # spill newly overlaps a neighbour declared box
    hard: bool = False  # autofit-off OR page-overflow OR collision (visually broken)


@dataclass
class SlideResult:
    page_index: int
    n_text_boxes: int = 0
    n_legal: int = 0
    hits: list = field(default_factory=list)


# --------------------------------------------------------------------------- #
# declared geometry from python-pptx
# --------------------------------------------------------------------------- #
def _auto_size_name(tf) -> str | None:
    try:
        a = tf.auto_size
        return None if a is None else str(a)
    except Exception:
        return None


def text_boxes_for_slide(slide, prefix: str) -> list[Box]:
    """Top-level shapes that carry visible text, as declared boxes (in points)."""
    boxes: list[Box] = []
    for i, shp in enumerate(slide.shapes):
        try:
            if not shp.has_text_frame:
                continue
        except Exception:
            continue
        text = (shp.text_frame.text or "").strip()
        if not text:
            continue
        if shp.left is None or shp.top is None or shp.width is None or shp.height is None:
            continue
        x0, y0 = emu2pt(shp.left), emu2pt(shp.top)
        boxes.append(
            Box(
                sid=f"{prefix}_e{i:02d}",
                x0=x0,
                y0=y0,
                x1=x0 + emu2pt(shp.width),
                y1=y0 + emu2pt(shp.height),
                text=text,
                auto_size=_auto_size_name(shp.text_frame),
                word_wrap=getattr(shp.text_frame, "word_wrap", None),
            )
        )
    return boxes


def _iou(a: Box, b: Box) -> float:
    ix0, iy0 = max(a.x0, b.x0), max(a.y0, b.y0)
    ix1, iy1 = min(a.x1, b.x1), min(a.y1, b.y1)
    iw, ih = max(0.0, ix1 - ix0), max(0.0, iy1 - iy0)
    inter = iw * ih
    if inter <= 0:
        return 0.0
    union = a.w * a.h + b.w * b.h - inter
    return inter / union if union > 0 else 0.0


def legal_boxes(
    boxes: list[Box], page_w: float, page_h: float, *, margin_pt: float, max_iou: float
) -> list[Box]:
    """G7 precondition: declared box in-margin AND non-overlapping with any other box."""
    legal: list[Box] = []
    for b in boxes:
        # in-margin: fully inside the page with a margin band on every side
        in_margin = (
            b.x0 >= margin_pt
            and b.y0 >= margin_pt
            and b.x1 <= page_w - margin_pt
            and b.y1 <= page_h - margin_pt
        )
        if not in_margin:
            continue
        overlaps = any(other is not b and _iou(b, other) > max_iou for other in boxes)
        if overlaps:
            continue
        legal.append(b)
    return legal


# --------------------------------------------------------------------------- #
# rendered text extent from the PDF
# --------------------------------------------------------------------------- #
def page_spans(page) -> list[tuple]:
    """(x0,y0,x1,y1,text) for every text span on a PyMuPDF page (points, TL origin)."""
    spans: list[tuple] = []
    d = page.get_text("dict")
    for blk in d.get("blocks", []):
        for line in blk.get("lines", []):
            for sp in line.get("spans", []):
                t = (sp.get("text") or "").strip()
                if not t:
                    continue
                x0, y0, x1, y1 = sp["bbox"]
                spans.append((x0, y0, x1, y1, t))
    return spans


def overflow_for_box(
    box: Box, spans: list[tuple], all_boxes: list[Box], page_w: float, page_h: float, *, tol: float
):
    """Rendered extent of this box's text; overflow past the declared box edges.

    Attribution exploits that we only score *legal* (non-overlapping) boxes: a box's
    text "owns" the region from its top edge down to the next declared box below it in
    the same horizontal band (and similarly to the right). Spilled wrapped lines that
    begin *below* the box bottom are therefore still attributed here, while the next
    box's own text is excluded by the neighbour bound — so the full spill magnitude is
    measured without stealing a neighbour's content.
    """
    # vertical ownership extends down to the nearest box below sharing this x-band
    below = [
        o.y0
        for o in all_boxes
        if o is not box and not (o.x1 <= box.x0 or o.x0 >= box.x1) and o.y0 >= box.y1 - tol
    ]
    v_limit = min(below) if below else page_h
    # horizontal ownership extends right to the nearest box sharing this y-band
    rights = [
        o.x0
        for o in all_boxes
        if o is not box and not (o.y1 <= box.y0 or o.y0 >= box.y1) and o.x0 >= box.x1 - tol
    ]
    h_limit = min(rights) if rights else page_w

    cx_lo, cx_hi = box.x0 - tol, box.x1 + tol
    cy_lo, cy_hi = box.y0 - tol, box.y1 + tol
    bottom = box.y1
    right = box.x1
    matched = False
    for x0, y0, x1, y1, _t in spans:
        cx = 0.5 * (x0 + x1)
        cy = 0.5 * (y0 + y1)
        # vertical: span horizontally within the box band, vertically owned by this box
        if cx_lo <= cx <= cx_hi and (box.y0 - tol) <= y0 < v_limit:
            bottom = max(bottom, y1)
            matched = True
        # horizontal: span vertically within the box band, horizontally owned by this box
        if cy_lo <= cy <= cy_hi and (box.x0 - tol) <= x0 < h_limit:
            right = max(right, x1)
            matched = True
    v_over = max(0.0, bottom - box.y1)
    h_over = max(0.0, right - box.x1)
    return matched, bottom, right, v_over, h_over


# --------------------------------------------------------------------------- #
# per-deck detection
# --------------------------------------------------------------------------- #
def detect_pptx(
    pptx_path: Path,
    *,
    margin_pt: float,
    max_iou: float,
    tol: float,
    min_over_pt: float,
    min_over_frac: float,
    workdir: Path,
) -> list[SlideResult]:
    import fitz  # PyMuPDF
    from pptx import Presentation

    from slide_examiner.render import render_pptx_to_pdf

    prs = Presentation(str(pptx_path))
    slide_w = emu2pt(prs.slide_width)
    slide_h = emu2pt(prs.slide_height)

    pdf_path = render_pptx_to_pdf(pptx_path, workdir)
    doc = fitz.open(str(pdf_path))

    results: list[SlideResult] = []
    for idx, slide in enumerate(prs.slides):
        if idx >= doc.page_count:
            break
        page = doc[idx]
        sx = page.rect.width / slide_w if slide_w else 1.0
        sy = page.rect.height / slide_h if slide_h else 1.0

        boxes = text_boxes_for_slide(slide, prefix=f"p{idx:02d}")
        # scale declared boxes into PDF point space (≈1.0 unless soffice rescales)
        for b in boxes:
            b.x0, b.x1, b.y0, b.y1 = b.x0 * sx, b.x1 * sx, b.y0 * sy, b.y1 * sy
        pg_w, pg_h = page.rect.width, page.rect.height

        legal = legal_boxes(boxes, pg_w, pg_h, margin_pt=margin_pt, max_iou=max_iou)
        spans = page_spans(page)

        sr = SlideResult(page_index=idx, n_text_boxes=len(boxes), n_legal=len(legal))
        for b in legal:
            matched, bottom, right, v_over, h_over = overflow_for_box(
                b, spans, boxes, pg_w, pg_h, tol=tol
            )
            if not matched:
                continue
            v_frac = v_over / b.h if b.h > 0 else 0.0
            h_frac = h_over / b.w if b.w > 0 else 0.0
            v_hit = v_over >= min_over_pt and v_frac >= min_over_frac
            h_hit = h_over >= min_over_pt and h_frac >= min_over_frac
            if v_hit or h_hit:
                # classify hardness: a hit is "hard" (visually broken, not gracefully
                # autofit-absorbed) if the box has autofit OFF, or the rendered content
                # spills off the page edge, or the spill newly collides with a neighbour.
                rend = Box(b.sid, b.x0, b.y0, max(right, b.x1), max(bottom, b.y1), "")
                collides = any(
                    o is not b and _iou(rend, o) > max_iou and _iou(b, o) <= max_iou
                    for o in boxes
                )
                page_over = max(0.0, max(bottom, b.y1) - pg_h, max(right, b.x1) - pg_w)
                autofit_off = (b.auto_size or "").startswith("NONE")
                hard = autofit_off or page_over >= min_over_pt or collides
                sr.hits.append(
                    Hit(
                        sid=b.sid,
                        box=(round(b.x0, 1), round(b.y0, 1), round(b.x1, 1), round(b.y1, 1)),
                        rendered_bottom=round(bottom, 1),
                        rendered_right=round(right, 1),
                        v_overflow_pt=round(v_over, 1),
                        h_overflow_pt=round(h_over, 1),
                        v_overflow_frac=round(v_frac, 3),
                        h_overflow_frac=round(h_frac, 3),
                        auto_size=b.auto_size,
                        text_preview=b.text[:60].replace("\n", " "),
                        page_overflow_pt=round(page_over, 1),
                        collides=collides,
                        hard=hard,
                    )
                )
        results.append(sr)
    doc.close()
    return results


# --------------------------------------------------------------------------- #
# CLI / aggregation
# --------------------------------------------------------------------------- #
def _iter_pptx(paths: list[str]) -> list[Path]:
    out: list[Path] = []
    for p in paths:
        pp = Path(p)
        if pp.is_dir():
            out.extend(sorted(pp.rglob("*.pptx")))
        elif pp.suffix.lower() == ".pptx":
            out.append(pp)
    # skip libreoffice lock/temp files
    return [p for p in out if not p.name.startswith(("~$", "."))]


def main() -> None:
    ap = argparse.ArgumentParser(description="G7 render-containment-overflow prevalence detector")
    ap.add_argument("paths", nargs="+", help="pptx files and/or directories")
    ap.add_argument("--margin-pt", type=float, default=18.0, help="page margin band (pt); 18pt=0.25in")
    ap.add_argument("--max-iou", type=float, default=0.05, help="declared-overlap threshold for 'legal'")
    ap.add_argument("--tol", type=float, default=2.0, help="attribution tolerance (pt)")
    ap.add_argument("--min-over-pt", type=float, default=3.0, help="min absolute overflow (pt) to flag")
    ap.add_argument("--min-over-frac", type=float, default=0.04, help="min relative overflow (of box dim)")
    ap.add_argument("--limit", type=int, default=0, help="cap number of decks (0=all)")
    ap.add_argument("--out", type=str, default="", help="write per-deck JSONL here")
    args = ap.parse_args()

    decks = _iter_pptx(args.paths)
    if args.limit:
        decks = decks[: args.limit]

    tot_decks = tot_slides = tot_text = tot_legal = tot_hits = 0
    decks_with_hit = slides_with_hit = 0
    tot_hard = decks_with_hard = slides_with_hard = 0
    errors = 0
    out_f = open(args.out, "w") if args.out else None

    with tempfile.TemporaryDirectory(prefix="g7prev-") as td:
        wd = Path(td)
        for di, deck in enumerate(decks):
            try:
                res = detect_pptx(
                    deck,
                    margin_pt=args.margin_pt,
                    max_iou=args.max_iou,
                    tol=args.tol,
                    min_over_pt=args.min_over_pt,
                    min_over_frac=args.min_over_frac,
                    workdir=wd,
                )
            except Exception as e:  # noqa: BLE001
                errors += 1
                print(f"  [ERR] {deck.name}: {type(e).__name__}: {e}", file=sys.stderr)
                continue
            d_text = sum(s.n_text_boxes for s in res)
            d_legal = sum(s.n_legal for s in res)
            d_hits = sum(len(s.hits) for s in res)
            d_slide_hit = sum(1 for s in res if s.hits)
            d_hard = sum(1 for s in res for h in s.hits if h.hard)
            d_slide_hard = sum(1 for s in res if any(h.hard for h in s.hits))
            tot_decks += 1
            tot_slides += len(res)
            tot_text += d_text
            tot_legal += d_legal
            tot_hits += d_hits
            slides_with_hit += d_slide_hit
            decks_with_hit += 1 if d_hits else 0
            tot_hard += d_hard
            slides_with_hard += d_slide_hard
            decks_with_hard += 1 if d_hard else 0
            if out_f:
                out_f.write(
                    json.dumps(
                        {
                            "deck": str(deck),
                            "slides": len(res),
                            "text_boxes": d_text,
                            "legal_boxes": d_legal,
                            "hits": d_hits,
                            "slides_with_hit": d_slide_hit,
                            "detail": [
                                {"page": s.page_index, "n_legal": s.n_legal, "hits": [vars(h) for h in s.hits]}
                                for s in res
                                if s.hits
                            ],
                        }
                    )
                    + "\n"
                )
            print(
                f"[{di+1}/{len(decks)}] {deck.name[:48]:48s} "
                f"slides={len(res):2d} legal={d_legal:3d} G7hits={d_hits:2d}"
                + ("  <<" if d_hits else "")
            )

    if out_f:
        out_f.close()

    def pct(a, b):
        return f"{100.0*a/b:.1f}%" if b else "n/a"

    print("\n================ G7 render-containment prevalence ================")
    print(f"decks scanned        : {tot_decks}  (errors: {errors})")
    print(f"slides               : {tot_slides}")
    print(f"text boxes           : {tot_text}")
    print(f"legal text boxes     : {tot_legal}  (G7 denominator)")
    print(f"G7 hits (boxes)      : {tot_hits}  (of which HARD: {tot_hard})")
    print("------------------------- ALL overflow (incl. soft autofit-grow) -")
    print(f"box-level incidence  : {pct(tot_hits, tot_legal)}   ({tot_hits}/{tot_legal} legal boxes)")
    print(f"slide-level incidence: {pct(slides_with_hit, tot_slides)}   ({slides_with_hit}/{tot_slides} slides)")
    print(f"deck-level incidence : {pct(decks_with_hit, tot_decks)}   ({decks_with_hit}/{tot_decks} decks)")
    print("------------------------- HARD G7 (autofit-off | off-page | collision) -")
    print(f"box-level incidence  : {pct(tot_hard, tot_legal)}   ({tot_hard}/{tot_legal} legal boxes)")
    print(f"slide-level incidence: {pct(slides_with_hard, tot_slides)}   ({slides_with_hard}/{tot_slides} slides)")
    print(f"deck-level incidence : {pct(decks_with_hard, tot_decks)}   ({decks_with_hard}/{tot_decks} decks)")
    print("=================================================================")
    print(
        f"params: margin={args.margin_pt}pt max_iou={args.max_iou} tol={args.tol}pt "
        f"min_over={args.min_over_pt}pt & {args.min_over_frac} frac"
    )


if __name__ == "__main__":
    main()
