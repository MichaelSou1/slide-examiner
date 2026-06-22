"""Part 3 R2 — extract a *lossless structured-geometry oracle* from real PPTX.

Closes the §8 hole: the cleanest discriminating experiment the paper could not
run was a *structured (image+oracle) evaluation on real slides*, because a real
oracle seemed to need human annotation (biased) or model self-annotation (also
biased). Zenodo10K ships 10k real CC-licensed ``.pptx`` whose element XML is
intact, so ``python-pptx`` extracts the geometry/text/style oracle **losslessly
from the source file** — no human, no model. That is exactly the perception
channel the attribution protocol needs, now on *real* layouts.

This module converts one PPTX slide into the repo's canonical ``Slide`` IR
(``slide_examiner.schemas.Slide``), in pixel space at 96 dpi (EMU/9525), aligned
to the VLM-SlideEval convention (PPTX XML GT + live render in one schema). The
same ``Slide`` then feeds modality B/C and the injector (``part3_real_inject``).

CLI (standalone dump of a deck sample):
  ~/anaconda3/envs/slide-examiner/bin/python scripts/part3_pptx_to_ir.py \
    --root data/part3/zenodo10k/pptx --out data/part3/real_ir --max-decks 30
"""
from __future__ import annotations

import argparse
import glob
import json
import os
from pathlib import Path

REPO = Path(__file__).resolve().parents[1]
import sys

sys.path.insert(0, str(REPO))

from slide_examiner.schemas import BBox, Element, Slide  # noqa: E402

EMU_PER_PX = 9525  # 914400 EMU/in / 96 px/in -> px @ 96 dpi (VLM-SlideEval convention)


# --------------------------------------------------------------------------- #
# python-pptx helpers (lazy import so the module loads without pptx installed)
# --------------------------------------------------------------------------- #
def _emu_to_px(v) -> float:
    return round(float(v) / EMU_PER_PX, 2) if v is not None else 0.0


def _shape_kind(shape) -> str:
    """Map a python-pptx shape to a coarse IR element type."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    try:
        st = shape.shape_type
    except Exception:  # noqa: BLE001
        st = None
    if st == MSO_SHAPE_TYPE.PICTURE:
        return "image"
    if getattr(shape, "has_table", False):
        return "table"
    if getattr(shape, "has_chart", False):
        return "chart"
    if st == MSO_SHAPE_TYPE.GROUP:
        return "group"
    # placeholder role -> title/subtitle/body
    try:
        if shape.is_placeholder:
            from pptx.enum.shapes import PP_PLACEHOLDER

            pt = shape.placeholder_format.type
            if pt in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                return "title"
            if pt == PP_PLACEHOLDER.SUBTITLE:
                return "subtitle"
            return "body"
    except Exception:  # noqa: BLE001
        pass
    if getattr(shape, "has_text_frame", False) and shape.has_text_frame and shape.text_frame.text.strip():
        return "body"
    return "shape"


def _placeholder_role(shape) -> str | None:
    try:
        if shape.is_placeholder:
            return str(shape.placeholder_format.type)
    except Exception:  # noqa: BLE001
        return None
    return None


def _first_font(shape):
    """Return (size_pt, color_hex) from the first run that declares them."""
    size_pt, color_hex = None, None
    if not getattr(shape, "has_text_frame", False) or not shape.has_text_frame:
        return size_pt, color_hex
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            f = run.font
            if size_pt is None and f.size is not None:
                try:
                    size_pt = round(f.size.pt, 1)
                except Exception:  # noqa: BLE001
                    pass
            if color_hex is None:
                try:
                    if f.color and f.color.type is not None and f.color.rgb is not None:
                        color_hex = "#" + str(f.color.rgb)
                except Exception:  # noqa: BLE001
                    pass
            if size_pt is not None and color_hex is not None:
                return size_pt, color_hex
        # fall back to paragraph-level font when runs carry none
        if size_pt is None and para.font.size is not None:
            try:
                size_pt = round(para.font.size.pt, 1)
            except Exception:  # noqa: BLE001
                pass
    return size_pt, color_hex


def _fill_hex(shape) -> str | None:
    try:
        fill = shape.fill
        if fill.type is not None:
            rgb = fill.fore_color.rgb
            if rgb is not None:
                return "#" + str(rgb)
    except Exception:  # noqa: BLE001
        return None
    return None


def _shape_text(shape) -> str:
    try:
        if getattr(shape, "has_table", False) and shape.has_table:
            cells = []
            for row in shape.table.rows:
                cells.append(" | ".join(c.text for c in row.cells))
            return "\n".join(cells).strip()
        if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
            return shape.text_frame.text.strip()
    except Exception:  # noqa: BLE001
        return ""
    return ""


def element_from_shape(shape, z: int) -> Element | None:
    """Convert one python-pptx shape to an IR Element (px @ 96 dpi)."""
    # connectors / shapes with no geometry are skipped
    try:
        left, top, width, height = shape.left, shape.top, shape.width, shape.height
    except Exception:  # noqa: BLE001
        return None
    if width is None or height is None or width <= 0 or height <= 0:
        return None
    size_pt, color_hex = _first_font(shape)
    style: dict = {}
    if size_pt is not None:
        style["font_size_pt"] = size_pt
    if color_hex is not None:
        style["color"] = color_hex
    fill = _fill_hex(shape)
    if fill is not None:
        style["fill_color"] = fill
    sid = f"shp{shape.shape_id}"
    meta: dict = {}
    role = _placeholder_role(shape)
    if role is not None:
        meta["role"] = role
    return Element(
        element_id=sid,
        type=_shape_kind(shape),
        bbox=BBox(_emu_to_px(left), _emu_to_px(top), _emu_to_px(width), _emu_to_px(height)),
        text=_shape_text(shape),
        style=style,
        z=z,
        placeholder_id=role,
        metadata=meta,
    )


def slide_to_ir(pptx_slide, *, slide_id: str, width_px: int, height_px: int) -> Slide:
    elements: list[Element] = []
    for z, shape in enumerate(pptx_slide.shapes):
        el = element_from_shape(shape, z)
        if el is not None:
            elements.append(el)
    return Slide(
        slide_id=slide_id,
        width=width_px,
        height=height_px,
        elements=tuple(elements),
        metadata={"scene": "full_proposal", "source": "zenodo10k"},
    )


def deck_dims_px(prs) -> tuple[int, int]:
    return int(round(prs.slide_width / EMU_PER_PX)), int(round(prs.slide_height / EMU_PER_PX))


# --------------------------------------------------------------------------- #
# Deck/slide sampling utilities (shared with the injector)
# --------------------------------------------------------------------------- #
def list_decks(root: str | Path, *, max_bytes: int = 25_000_000) -> list[Path]:
    """All .pptx under root, skipping oversized decks (load-twice + render cost)."""
    paths = [Path(p) for p in glob.glob(str(Path(root) / "**" / "*.pptx"), recursive=True)]
    out = [p for p in paths if p.stat().st_size <= max_bytes]
    return sorted(out, key=lambda p: p.stat().st_size)


def license_of(path: Path, root: Path) -> str:
    try:
        rel = path.relative_to(root)
        return rel.parts[0]
    except Exception:  # noqa: BLE001
        return "unknown"


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--root", default="data/part3/zenodo10k/pptx")
    ap.add_argument("--out", default="data/part3/real_ir")
    ap.add_argument("--max-decks", type=int, default=30)
    ap.add_argument("--max-bytes", type=int, default=25_000_000)
    args = ap.parse_args()

    from pptx import Presentation

    root = REPO / args.root if not Path(args.root).is_absolute() else Path(args.root)
    out_dir = REPO / args.out if not Path(args.out).is_absolute() else Path(args.out)
    out_dir.mkdir(parents=True, exist_ok=True)

    decks = list_decks(root, max_bytes=args.max_bytes)[: args.max_decks]
    index = []
    n_slides = 0
    for deck_path in decks:
        try:
            prs = Presentation(str(deck_path))
        except Exception as exc:  # noqa: BLE001
            print(f"  SKIP {deck_path.name}: {exc}")
            continue
        w, h = deck_dims_px(prs)
        deck_key = deck_path.stem[:48].replace(" ", "_")
        for si, slide in enumerate(prs.slides):
            ir = slide_to_ir(slide, slide_id=f"{deck_key}__s{si}", width_px=w, height_px=h)
            if not ir.elements:
                continue
            rel = str(deck_path.relative_to(root)) if str(root) in str(deck_path) else str(deck_path)
            fn = out_dir / f"{ir.slide_id}.json"
            fn.write_text(json.dumps(ir.to_dict(), ensure_ascii=False, indent=2), encoding="utf-8")
            index.append({
                "slide_id": ir.slide_id, "ir_path": str(fn), "deck_path": str(deck_path),
                "deck_rel": rel, "slide_index": si, "n_elements": len(ir.elements),
                "width": w, "height": h, "license": license_of(deck_path, root),
            })
            n_slides += 1
    (out_dir / "index.json").write_text(json.dumps(index, ensure_ascii=False, indent=2), encoding="utf-8")
    print(f"[done] {len(decks)} decks -> {n_slides} slide IRs -> {out_dir}")
    by_lic: dict = {}
    for r in index:
        by_lic[r["license"]] = by_lic.get(r["license"], 0) + 1
    print(f"  by license: {by_lic}")


if __name__ == "__main__":
    main()
