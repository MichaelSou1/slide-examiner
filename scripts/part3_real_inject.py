"""Part 3 R2 — inject taxonomy defects into REAL slide layouts (PPTX-native).

For each sampled real slide we build a one-slide ``.pptx`` (keep only the target
slide so LibreOffice renders exactly one page), then a *defective* copy whose
single shape is mutated in PPTX XML space (move / resize / recolor / overstuff).
BOTH the clean and defective one-slide decks are rendered by the SAME real
renderer (LibreOffice headless -> PDF -> PNG), so the paired image diff isolates
exactly the injected defect on *real pixels*. The lossless oracle (modality B/C)
is re-extracted from each PPTX with ``part3_pptx_to_ir`` — no human or model
annotation, so the §8 self-annotation-bias concern does not arise.

Classes injected (the ones that map cleanly to a single-shape PPTX mutation and
that carry the perception/capability story):
  G1 text overflow      autofit off + enlarge font  -> renders spill (perception/format)
  G2 element overlap     move a shape onto a sibling -> IR shows overlap (perception)
  G3 alignment offset    nudge a shape off its column -> IR shows offset (sub-perceptual)
  G4 font-size inconsist. resize one text block vs peers -> IR shows mismatch
  G6 margin violation    push a shape past the slide edge -> IR shows bleed
G5 (brand colour) and G7 (render containment) are out of real-deck scope: G5
needs a brand palette third-party decks do not declare; G7 is already the
synthetic falsifiable class (legal IR, overflow only in the HTML builder) and
cannot be reproduced from a real text frame without leaking the overflow text
into the oracle. The S-* semantic classes need content understanding to inject
faithfully and stay synthetic. We report this scope honestly.

A self-check drops any pair whose injection did NOT change the rendered pixels
(the §4 "snap/absorption" hazard, measured here on real free-form decks).

Usage:
  SLIDE_EXAMINER_SOFFICE=/tmp/lo_root/opt/libreoffice26.2/program/soffice \
  ~/anaconda3/envs/slide-examiner/bin/python scripts/part3_real_inject.py \
    --index data/part3/real_ir/index.json --per-class 40 \
    --out data/part3/manifest_real_rendered.jsonl --work data/part3/real_build
"""
from __future__ import annotations

import argparse
import json
import os
import random
import shutil
import subprocess
import statistics
from pathlib import Path

REPO = Path(__file__).resolve().parents[1]
import sys

sys.path.insert(0, str(REPO))
sys.path.insert(0, str(REPO / "scripts"))

import numpy as np  # noqa: E402
from PIL import Image  # noqa: E402

from slide_examiner.render import build_render_spec, image_size  # noqa: E402
from part3_pptx_to_ir import EMU_PER_PX, deck_dims_px, slide_to_ir  # noqa: E402

CLASSES = ["G1_TEXT_OVERFLOW", "G2_ELEMENT_OVERLAP", "G3_ALIGNMENT_OFFSET",
           "G4_FONT_SIZE_INCONSISTENCY", "G6_MARGIN_VIOLATION"]
CHANGE_THR = 20      # per-pixel channel-sum delta (same as part3_p3_fidelity)
CHANGE_EPS = 0.0008  # changed-pixel fraction above which the perturbation "rendered"


# --------------------------------------------------------------------------- #
# PPTX surgery
# --------------------------------------------------------------------------- #
def keep_only(prs, idx: int) -> None:
    """Drop every slide except ``idx`` from the sldIdLst so the deck renders to
    a single page (parts stay in the package; they are just unreferenced)."""
    sldIdLst = prs.slides._sldIdLst
    ids = list(sldIdLst)
    for i, sld in enumerate(ids):
        if i != idx:
            sldIdLst.remove(sld)


def _geom_shapes(slide):
    """Shapes with explicit, positive geometry (px), as (shape, x, y, w, h)."""
    out = []
    for sh in slide.shapes:
        try:
            l, t, w, h = sh.left, sh.top, sh.width, sh.height
        except Exception:  # noqa: BLE001
            continue
        if None in (l, t, w, h) or w <= 0 or h <= 0:
            continue
        out.append((sh, l / EMU_PER_PX, t / EMU_PER_PX, w / EMU_PER_PX, h / EMU_PER_PX))
    return out


def _text_shapes(slide, min_chars=1):
    out = []
    for sh, x, y, w, h in _geom_shapes(slide):
        if getattr(sh, "has_text_frame", False) and sh.has_text_frame:
            runs = [r for p in sh.text_frame.paragraphs for r in p.runs]
            if runs and len(sh.text_frame.text.strip()) >= min_chars:
                out.append((sh, x, y, w, h, runs))
    return out


def _median_font(slide) -> float:
    sizes = []
    for sh in slide.shapes:
        if getattr(sh, "has_text_frame", False) and sh.has_text_frame:
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    if r.font.size is not None:
                        sizes.append(r.font.size.pt)
    return statistics.median(sizes) if sizes else 18.0


# --------------------------------------------------------------------------- #
# Mutators — each returns (target_shape_id:int, severity:float, meta:dict) | None
# --------------------------------------------------------------------------- #
def _pin(sh, x, y, w, h):
    """Write ALL four geometry coords explicitly (px -> EMU). A placeholder that
    inherits its position has no <a:off>/<a:ext>; setting just one coordinate makes
    python-pptx synthesize the xfrm with the OTHERS defaulting to 0 (silently
    teleporting the shape to the top-left). Pinning all four keeps the injected
    perturbation 1-D and controlled."""
    from pptx.util import Emu
    sh.left = Emu(int(round(x * EMU_PER_PX)))
    sh.top = Emu(int(round(y * EMU_PER_PX)))
    sh.width = Emu(int(round(w * EMU_PER_PX)))
    sh.height = Emu(int(round(h * EMU_PER_PX)))


def _aligned_column(slide, tol=12.0, min_members=3):
    """>=min_members shapes that share (within ``tol`` px) a left edge — a visible
    alignment column (e.g. a bullet/text-box stack). The basis for an INTERNAL-contrast
    G3 (E8 re-operationalisation): one member nudged out of the column is decidable from
    the slide alone, no invisible external expected-position. Returns the members sorted
    top-to-bottom, or None when the real layout has no >=3-deep column."""
    buckets: dict[int, list] = {}
    for c in _geom_shapes(slide):
        buckets.setdefault(round(c[1] / tol), []).append(c)
    groups = [g for g in buckets.values() if len(g) >= min_members]
    if not groups:
        return None
    group = max(groups, key=len)  # the deepest column
    group.sort(key=lambda c: c[2])  # by y, top-to-bottom
    return group


def mut_g3(slide, W, H, rng):
    """INTERNAL contrast: shift ONE member of an aligned left-edge column out of line
    (well-posed — the outlier is decidable against its visible siblings)."""
    group = _aligned_column(slide)
    if group is None:
        return None
    sh, x, y, w, h = group[len(group) // 2]  # the middle member — clearly out of column
    offset = rng.choice([18, 26, 34, 44])
    nx = x + offset
    if nx + w > W - 2:           # would bleed off-slide -> shift the other way instead
        nx = max(2.0, x - offset)
    _pin(sh, nx, y, w, h)
    col_left = sum(c[1] for c in group) / len(group)
    return sh.shape_id, float(offset), {"offset_px": offset, "axis": "x", "mode": "internal",
                                        "column_left_px": round(col_left, 1), "column_n": len(group)}


def mut_g2(slide, W, H, rng):
    cands = _geom_shapes(slide)
    if len(cands) < 2:
        return None
    cands.sort(key=lambda c: c[3] * c[4], reverse=True)
    base = cands[0]
    movers = [c for c in cands[1:] if c[3] < base[3] * 1.2]  # smaller-ish movers
    if not movers:
        movers = cands[1:]
    sh, x, y, w, h = rng.choice(movers)
    bx, by, bw, bh = base[1], base[2], base[3], base[4]
    nx = max(2, min(bx + bw * 0.30, W - w - 2))
    ny = max(2, min(by + bh * 0.30, H - h - 2))
    _pin(sh, nx, ny, w, h)
    return sh.shape_id, 0.3, {"iou_target": 0.3, "base_id": f"shp{base[0].shape_id}"}


def mut_g6(slide, W, H, rng):
    cands = _geom_shapes(slide)
    if not cands:
        return None
    sh, x, y, w, h = rng.choice(cands)
    bleed = rng.choice([40, 60, 90])
    # push toward the nearest horizontal edge, keeping the vertical position fixed
    if x + w / 2 < W / 2:
        side, nx = "left", -bleed
    else:
        side, nx = "right", (W - w + bleed)
    _pin(sh, nx, y, w, h)
    return sh.shape_id, float(bleed), {"bleed_px": bleed, "side": side}


def mut_g4(slide, W, H, rng):
    from pptx.util import Pt
    texts = _text_shapes(slide, min_chars=3)
    if not texts:
        return None
    med = _median_font(slide)
    sh, x, y, w, h, runs = rng.choice(texts)
    cur = next((r.font.size.pt for r in runs if r.font.size is not None), med)
    if cur <= med + 1:  # enlarge
        new = max(med * 1.9, cur * 1.9)
    else:               # shrink an already-large block
        new = max(8.0, min(med * 0.5, cur * 0.5))
    for r in runs:
        r.font.size = Pt(round(new, 1))
    return sh.shape_id, abs(new - med), {"delta_pt": round(abs(new - med), 1),
                                         "new_pt": round(new, 1), "peer_pt": round(med, 1)}


def mut_g1(slide, W, H, rng):
    from pptx.util import Pt
    from pptx.enum.text import MSO_AUTO_SIZE
    texts = [t for t in _text_shapes(slide, min_chars=25)]
    if not texts:
        return None
    sh, x, y, w, h, runs = rng.choice(texts)
    tf = sh.text_frame
    try:
        tf.word_wrap = False
        tf.auto_size = MSO_AUTO_SIZE.NONE
    except Exception:  # noqa: BLE001
        pass
    cur = next((r.font.size.pt for r in runs if r.font.size is not None), _median_font(slide))
    new = cur * rng.choice([1.7, 2.0, 2.3])
    for r in runs:
        r.font.size = Pt(round(new, 1))
    return sh.shape_id, float(round(new - cur, 1)), {"from_pt": round(cur, 1), "to_pt": round(new, 1),
                                                     "mode": "enlarge_nowrap"}


MUTATORS = {"G1_TEXT_OVERFLOW": mut_g1, "G2_ELEMENT_OVERLAP": mut_g2,
            "G3_ALIGNMENT_OFFSET": mut_g3, "G4_FONT_SIZE_INCONSISTENCY": mut_g4,
            "G6_MARGIN_VIOLATION": mut_g6}


# --------------------------------------------------------------------------- #
# Rendering (batched LibreOffice)
# --------------------------------------------------------------------------- #
def _soffice() -> str:
    s = os.environ.get("SLIDE_EXAMINER_SOFFICE")
    if s and Path(s).exists():
        return s
    return shutil.which("soffice") or shutil.which("libreoffice")


def batch_render(pptx_paths: list[Path], workdir: Path, dpi: int = 150,
                 chunk: int = 40) -> dict[str, str]:
    """Convert many one-slide pptx -> one png each, in chunked soffice calls."""
    soffice = _soffice()
    if not soffice:
        raise RuntimeError("no soffice; set SLIDE_EXAMINER_SOFFICE")
    pdf_dir = workdir / "pdf"
    png_dir = workdir / "png"
    pdf_dir.mkdir(parents=True, exist_ok=True)
    png_dir.mkdir(parents=True, exist_ok=True)
    prof = workdir / "lo_prof"
    out: dict[str, str] = {}
    for i in range(0, len(pptx_paths), chunk):
        batch = pptx_paths[i:i + chunk]
        shutil.rmtree(prof, ignore_errors=True)
        prof.mkdir(parents=True, exist_ok=True)
        env = {**os.environ, "HOME": str(prof)}
        cmd = [soffice, "--headless", "--invisible", "--nodefault", "--nolockcheck",
               "--nologo", "--norestore", f"-env:UserInstallation=file://{prof}",
               "--convert-to", "pdf", "--outdir", str(pdf_dir)] + [str(p) for p in batch]
        subprocess.run(cmd, env=env, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL,
                       stdin=subprocess.DEVNULL, timeout=900, check=False)
        for p in batch:
            pdf = pdf_dir / (p.stem + ".pdf")
            if not pdf.exists():
                continue
            base = png_dir / p.stem
            subprocess.run(["pdftoppm", "-png", "-r", str(dpi), "-f", "1", "-l", "1",
                            str(pdf), str(base)], stdout=subprocess.DEVNULL,
                           stderr=subprocess.DEVNULL, stdin=subprocess.DEVNULL, check=False)
            # pdftoppm pads the page number by total pages (1 page -> "-1")
            for cand in (Path(f"{base}-1.png"), Path(f"{base}-01.png"), Path(f"{base}.png")):
                if cand.exists():
                    out[str(p)] = str(cand)
                    break
        print(f"  [render] {min(i + chunk, len(pptx_paths))}/{len(pptx_paths)}")
    return out


def changed_frac(a_path: str, b_path: str) -> float | None:
    try:
        a = Image.open(a_path).convert("RGB")
        b = Image.open(b_path).convert("RGB").resize(a.size)
    except Exception:  # noqa: BLE001
        return None
    aa = np.asarray(a, dtype=np.int16)
    bb = np.asarray(b, dtype=np.int16)
    return float((np.abs(aa - bb).sum(-1) > CHANGE_THR).mean())


# --------------------------------------------------------------------------- #
# Build
# --------------------------------------------------------------------------- #
def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--index", default="data/part3/real_ir/index.json")
    ap.add_argument("--out", default="data/part3/manifest_real_rendered.jsonl")
    ap.add_argument("--fidelity-out", default="data/part3/pc_real_fidelity.json")
    ap.add_argument("--work", default="data/part3/real_build")
    ap.add_argument("--per-class", type=int, default=40)
    ap.add_argument("--min-elements", type=int, default=2)
    ap.add_argument("--max-per-deck", type=int, default=14)
    ap.add_argument("--dpi", type=int, default=150)
    ap.add_argument("--seed", type=int, default=20260622)
    ap.add_argument("--classes", nargs="+", default=None,
                    help="restrict to a subset of CLASSES (e.g. G3_ALIGNMENT_OFFSET for the "
                         "E8 internal-G3 regen) so other classes' frozen renders are untouched.")
    args = ap.parse_args()

    global CLASSES
    if args.classes:
        CLASSES = [c for c in CLASSES if c in set(args.classes)]
        if not CLASSES:
            raise SystemExit(f"--classes {args.classes} matched none of {list(MUTATORS)}")

    from pptx import Presentation

    idx = json.loads((REPO / args.index).read_text())
    work = REPO / args.work
    pptx_dir = work / "pptx"
    ir_dir = work / "ir"
    for d in (pptx_dir, ir_dir):
        d.mkdir(parents=True, exist_ok=True)
    rng = random.Random(args.seed)

    # candidate slides: enough elements, shuffled, capped per deck
    cands = [r for r in idx if r["n_elements"] >= args.min_elements]
    rng.shuffle(cands)
    per_deck: dict[str, int] = {}
    pool = []
    for r in cands:
        dk = r["deck_path"]
        if per_deck.get(dk, 0) >= args.max_per_deck:
            continue
        per_deck[dk] = per_deck.get(dk, 0) + 1
        pool.append(r)

    # round-robin assign classes to slides; build one-slide clean+defective pptx
    plan = []  # dicts with paths + meta
    counts = {c: 0 for c in CLASSES}
    pi = 0
    target_total = args.per_class * len(CLASSES)
    for r in pool:
        if sum(counts.values()) >= target_total:
            break
        # pick a class still under quota (rotate to balance)
        order = sorted(CLASSES, key=lambda c: counts[c])
        made = False
        for cls in order:
            if counts[cls] >= args.per_class:
                continue
            try:
                prs = Presentation(r["deck_path"])
            except Exception:  # noqa: BLE001
                break
            si = r["slide_index"]
            if si >= len(prs.slides._sldIdLst):
                break
            keep_only(prs, si)
            W, H = deck_dims_px(prs)
            stem = f"{r['slide_id']}__{cls}"
            clean_pptx = pptx_dir / f"{stem}__clean.pptx"
            prs.save(str(clean_pptx))
            clean_ir = slide_to_ir(prs.slides[0], slide_id=stem, width_px=W, height_px=H)
            # defective: load the FULL deck fresh (NOT a reopen of the saved one-slide
            # deck — reopening + re-saving compounds python-pptx's duplicate-part bug
            # and corrupts the package so LibreOffice silently drops it).
            prs2 = Presentation(r["deck_path"])
            keep_only(prs2, si)
            res = MUTATORS[cls](prs2.slides[0], W, H, rng)
            if res is None:
                clean_pptx.unlink(missing_ok=True)
                continue
            tgt_id, severity, meta = res
            def_pptx = pptx_dir / f"{stem}__def.pptx"
            prs2.save(str(def_pptx))
            def_ir = slide_to_ir(prs2.slides[0], slide_id=stem, width_px=W, height_px=H)
            clean_ir_path = ir_dir / f"{stem}__clean.json"
            def_ir_path = ir_dir / f"{stem}__def.json"
            clean_ir_path.write_text(json.dumps(clean_ir.to_dict(), ensure_ascii=False), encoding="utf-8")
            def_ir_path.write_text(json.dumps(def_ir.to_dict(), ensure_ascii=False), encoding="utf-8")
            plan.append({
                "sample_id": stem, "defect": cls, "target_id": f"shp{tgt_id}",
                "severity": severity, "label_meta": meta, "W": W, "H": H,
                "deck_path": r["deck_path"], "slide_index": si, "license": r.get("license"),
                "clean_pptx": clean_pptx, "def_pptx": def_pptx,
                "clean_ir": clean_ir, "def_ir": def_ir,
                "clean_ir_path": str(clean_ir_path), "def_ir_path": str(def_ir_path),
            })
            counts[cls] += 1
            made = True
            break
        pi += 1
    print(f"[plan] {len(plan)} instances; per-class counts: {counts}")

    # batch render every clean + defective one-slide pptx
    all_pptx = [p["clean_pptx"] for p in plan] + [p["def_pptx"] for p in plan]
    rendered = batch_render(all_pptx, work, dpi=args.dpi)

    # assemble manifest + absorption self-check
    records, absorbed, no_render, ir_unfaithful = [], 0, 0, 0
    per_class_kept = {c: 0 for c in CLASSES}
    fidelity_rows = []  # every planned instance, kept or not (for the render-fidelity audit)
    for p in plan:
        cpng = rendered.get(str(p["clean_pptx"]))
        dpng = rendered.get(str(p["def_pptx"]))
        if not cpng or not dpng:
            no_render += 1
            fidelity_rows.append({"defect": p["defect"], "changed_frac": None, "status": "no_render"})
            continue
        cf = changed_frac(dpng, cpng)
        if cf is None:
            no_render += 1
            fidelity_rows.append({"defect": p["defect"], "changed_frac": None, "status": "no_render"})
            continue
        rendered_ok = cf > CHANGE_EPS
        fidelity_rows.append({"defect": p["defect"], "changed_frac": round(cf, 5),
                              "status": "rendered" if rendered_ok else "absorbed"})
        if not rendered_ok:
            absorbed += 1
            continue  # injection did not change the real pixels -> drop (absorbed)
        # IR-faithfulness self-check (G3 internal contrast): the offset must survive into
        # the lossless oracle so the structure channel (modality B/C) is well-posed. A few
        # placeholder-inherited shapes reflow in pixels but barely move in the IR -> drop.
        if p["defect"] == "G3_ALIGNMENT_OFFSET":
            off = float(p["label_meta"].get("offset_px", 0.0))
            cx = {e["element_id"]: e for e in p["clean_ir"].to_dict()["elements"]}
            dx = {e["element_id"]: e for e in p["def_ir"].to_dict()["elements"]}
            ce, de = cx.get(p["target_id"]), dx.get(p["target_id"])
            ir_dx = abs(de["bbox"]["x"] - ce["bbox"]["x"]) if ce and de else 0.0
            if ir_dx < off * 0.5:
                fidelity_rows[-1]["status"] = "ir_unfaithful"
                ir_unfaithful += 1  # rendered, but defect not legible in the oracle -> drop
                continue
        iw, ih = image_size(dpng)
        spec = build_render_spec(slide_width=p["W"], slide_height=p["H"], image_width=iw,
                                 image_height=ih, renderer="libreoffice", dpi=args.dpi)
        ciw, cih = image_size(cpng)
        cspec = build_render_spec(slide_width=p["W"], slide_height=p["H"], image_width=ciw,
                                  image_height=cih, renderer="libreoffice", dpi=args.dpi)
        records.append({
            "sample_id": p["sample_id"], "defect": p["defect"],
            "image_path": str(Path(dpng).resolve()),
            "slide": p["def_ir"].to_dict(),
            "clean_slide": p["clean_ir"].to_dict(),
            "labels": [{"type": p["defect"], "severity": p["severity"],
                        "target_element_ids": [p["target_id"]], "metadata": p["label_meta"]}],
            "pair": {
                "clean_image_path": str(Path(cpng).resolve()),
                "defective_image_path": str(Path(dpng).resolve()),
                "clean_slide_path": p["clean_ir_path"],
                "defective_slide_path": p["def_ir_path"],
            },
            "metadata": {
                "scene": "full_proposal", "source": "zenodo10k", "deck_path": p["deck_path"],
                "slide_index": p["slide_index"], "license": p["license"],
                "render": spec, "render_clean": cspec, "changed_frac": round(cf, 5),
                "clean_image_path": str(Path(cpng).resolve()),
                "defective_image_path": str(Path(dpng).resolve()),
            },
        })
        per_class_kept[p["defect"]] += 1

    out = REPO / args.out
    out.write_text("".join(json.dumps(r, ensure_ascii=False) + "\n" for r in records), encoding="utf-8")

    # render-fidelity audit (real free-form decks vs the synthetic snap-to-master 45% hazard)
    n = len(plan)
    fid: dict = {"n_injected": n, "thr_channel_sum": CHANGE_THR, "eps_changed_frac": CHANGE_EPS,
                 "per_defect": {}, "overall": {}}
    rendered_tot = 0
    for cls in CLASSES:
        cr = [r for r in fidelity_rows if r["defect"] == cls]
        scored = [r for r in cr if r["changed_frac"] is not None]
        rend = [r for r in scored if r["status"] == "rendered"]
        rendered_tot += len(rend)
        cfs = sorted(r["changed_frac"] for r in scored)
        fid["per_defect"][cls] = {
            "n_injected": len(cr),
            "ir_injected_rate": 1.0,  # injection is in IR by construction
            "rendered_rate": round(len(rend) / len(scored), 3) if scored else None,
            "absorption_rate": round(sum(r["status"] == "absorbed" for r in scored) / len(scored), 3) if scored else None,
            "changed_frac_median": round(cfs[len(cfs) // 2], 5) if cfs else None,
        }
    n_scored = sum(1 for r in fidelity_rows if r["changed_frac"] is not None)
    fid["overall"] = {
        "n_scored": n_scored, "rendered_rate": round(rendered_tot / n_scored, 3) if n_scored else None,
        "absorption_rate": round(absorbed / n_scored, 3) if n_scored else None,
        "ir_unfaithful": ir_unfaithful,  # G3: rendered but offset not legible in the oracle (dropped)
        "note": ("Real free-form decks are NOT snap-to-master: injected geometry defects render "
                 "faithfully (rendered_rate near 1.0), unlike the synthetic enterprise template that "
                 "silently absorbs 45% of injected geometry defects (paper Sec.4 / Result-3b)."),
    }
    fid_out = REPO / args.fidelity_out
    fid_out.write_text(json.dumps(fid, indent=2, ensure_ascii=False), encoding="utf-8")

    print(f"[done] kept {len(records)}/{n} pairs -> {out}")
    print(f"  per-class kept: {per_class_kept}")
    print(f"  absorbed (rendered identical) = {absorbed}/{n} = {absorbed / n:.3f}" if n else "")
    print(f"  ir-unfaithful (rendered but offset not in oracle) = {ir_unfaithful}/{n}")
    print(f"  no-render (soffice/pdftoppm miss) = {no_render}/{n}")
    print(f"  render-fidelity -> {fid_out}: overall rendered_rate={fid['overall']['rendered_rate']}")


if __name__ == "__main__":
    main()
